# services/ppt_renderer.py
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import shutil

# ==================== 母版背景设置函数 ====================
def set_master_background(prs, bg_config):
    """
    设置母版背景（替换所有幻灯片的默认背景）
    bg_config: 背景配置字典
    """
    if not bg_config:
        return False
    
    bg_type = bg_config.get("type", "solid")
    
    try:
        # 获取母版
        master = prs.slide_master
        
        if bg_type == "solid":
            # 纯色背景
            color = bg_config.get("color", [255, 255, 255])
            master.background.fill.solid()
            master.background.fill.fore_color.rgb = RGBColor(color[0], color[1], color[2])
            print(f"✅ 母版纯色背景设置成功: RGB{color}")
            
        elif bg_type == "gradient":
            # 渐变背景
            color1 = bg_config.get("color", [200, 220, 255])
            color2 = bg_config.get("color2", [100, 150, 255])
            
            master.background.fill.gradient()
            master.background.fill.gradient_angle = 45
            
            stops = master.background.fill.gradient_stops
            stops[0].color.rgb = RGBColor(color1[0], color1[1], color1[2])
            stops[0].position = 0.0
            stops[1].color.rgb = RGBColor(color2[0], color2[1], color2[2])
            stops[1].position = 1.0
            
            print(f"✅ 母版渐变背景设置成功")
        
        return True
        
    except Exception as e:
        print(f"❌ 母版背景设置失败: {e}")
        return False

# ==================== 版式0：比较页 ====================
def fill_compare_slide(slide, data):
    """填充比较页 (版式0)
    占位符: idx=0(标题), idx=1,2,3,4(四个内容区), idx=10,11,12(页脚)
    """
    try:
        slide.placeholders[0].text = data.get("title", "")
        
        # 左上 (idx=1)
        if "left_title" in data:
            slide.placeholders[1].text = data["left_title"]
        
        # 左下 (idx=2)
        if "left_items" in data:
            tf = slide.placeholders[2].text_frame
            tf.clear()
            for i, item in enumerate(data["left_items"]):
                if i == 0:
                    tf.paragraphs[0].text = f"• {item}"
                else:
                    p = tf.add_paragraph()
                    p.text = f"• {item}"
        
        # 右上 (idx=3)
        if "right_title" in data:
            slide.placeholders[3].text = data["right_title"]
        
        # 右下 (idx=4)
        if "right_items" in data:
            tf = slide.placeholders[4].text_frame
            tf.clear()
            for i, item in enumerate(data["right_items"]):
                if i == 0:
                    tf.paragraphs[0].text = f"• {item}"
                else:
                    p = tf.add_paragraph()
                    p.text = f"• {item}"
        
        print("    ✓ 比较页填充成功")
    except Exception as e:
        print(f"    ✗ 比较页填充失败: {e}")

# ==================== 版式1：末尾幻灯片 ====================
def fill_end_slide(slide, data):
    """填充末尾幻灯片 (版式1)
    占位符: idx=10,11,12(页脚), idx=0(标题), idx=13(内容)
    """
    try:
        # 标题 (idx=0)
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:
                placeholder.text = data.get("title", "谢谢观看")
                break
        
        # 内容 (idx=13)
        if "content" in data:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 13:
                    placeholder.text = data["content"]
                    break
        
        print("    ✓ 末尾页填充成功")
    except Exception as e:
        print(f"    ✗ 末尾页填充失败: {e}")

# ==================== 版式2：我的开头幻灯片 ====================
def fill_my_title_slide(slide, data):
    """填充自定义标题页 (版式2)
    占位符: idx=0(标题), idx=10,11,12(页脚), idx=13(副标题)
    """
    try:
        # 标题 (idx=0)
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:
                placeholder.text = data.get("title", "")
                break
        
        # 副标题 (idx=13)
        if "subtitle" in data:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 13:
                    placeholder.text = data["subtitle"]
                    break
        
        print("    ✓ 自定义标题页填充成功")
    except Exception as e:
        print(f"    ✗ 自定义标题页填充失败: {e}")

# ==================== 版式3：表格页 ====================
def fill_table_slide(slide, data):
    """填充表格页 (版式3) - 完全自适应版
    占位符: idx=0(标题), idx=10,11,12(页脚), idx=14(表格)
    """
    try:
        # 标题 (idx=0)
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:
                placeholder.text = data.get("title", "")
                for paragraph in placeholder.text_frame.paragraphs:
                    paragraph.font.size = Pt(28)
                    paragraph.font.bold = True
                break
        
        # 获取表格数据
        table_data = data.get("table_data", [])
        if not table_data:
            print("    ⚠️ 没有表格数据")
            return
        
        rows = len(table_data)
        cols = len(table_data[0])
        
        # 查找表格占位符
        table_placeholder = None
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 14:
                table_placeholder = placeholder
                break
        
        if not table_placeholder:
            print("    ❌ 未找到表格占位符 (idx=14)")
            return
        
        # 获取占位符大小
        placeholder_width = table_placeholder.width
        placeholder_height = table_placeholder.height
        
        # ===== 1. 自适应行高计算 =====
        available_height = placeholder_height - Inches(0.4)
        row_height = available_height / rows
        
        min_row_height = Inches(0.3)
        if row_height < min_row_height:
            row_height = min_row_height
            total_height = row_height * rows + Inches(0.4)
        else:
            total_height = placeholder_height
        
        # ===== 2. 自适应字体大小计算函数 =====
        def calc_font_size(text, cell_width, cell_height, is_header=False):
            base_size = 24 if is_header else 18
            text_len = len(str(text))
            if text_len > 20:
                base_size = max(12, base_size - 6)
            elif text_len > 10:
                base_size = max(14, base_size - 4)
            elif text_len > 5:
                base_size = max(16, base_size - 2)
            
            cell_height_inches = cell_height / 914400
            max_height_size = int(cell_height_inches * 18)
            base_size = min(base_size, max_height_size)
            
            min_size = 12 if is_header else 10
            base_size = max(base_size, min_size)
            return Pt(base_size)
        
        # ===== 3. 计算最佳列宽 =====
        col_max_lengths = [0] * cols
        for i in range(rows):
            for j in range(cols):
                text_len = len(str(table_data[i][j]))
                col_max_lengths[j] = max(col_max_lengths[j], text_len)
        
        total_chars = sum(col_max_lengths)
        if total_chars > 0:
            available_width = placeholder_width - Inches(0.2 * cols)
            col_widths = []
            for length in col_max_lengths:
                ratio = length / total_chars
                width = int(available_width * ratio)
                min_width = int(Inches(1.0))
                col_widths.append(max(width, min_width))
        else:
            col_width = int(placeholder_width / cols)
            col_widths = [col_width] * cols
        
        # 删除原占位符
        sp = table_placeholder._element
        sp.getparent().remove(sp)
        
        # 创建新表格
        left = table_placeholder.left
        top = table_placeholder.top
        
        table_shape = slide.shapes.add_table(
            rows, cols, 
            left, top, 
            placeholder_width, total_height
        )
        table = table_shape.table
        
        # 设置行高
        for i in range(rows):
            table.rows[i].height = int(row_height)
        
        # 设置列宽
        for j, width in enumerate(col_widths):
            table.columns[j].width = width
        
        # ===== 4. 填充数据 =====
        for i in range(rows):
            for j in range(cols):
                cell = table.cell(i, j)
                value = table_data[i][j]
                cell.text = str(value)
                
                cell.margin_left = Inches(0.08)
                cell.margin_right = Inches(0.08)
                cell.margin_top = Inches(0.04)
                cell.margin_bottom = Inches(0.04)
                
                cell_width = table.columns[j].width
                cell_height = table.rows[i].height
                
                is_header = (i == 0)
                font_size = calc_font_size(value, cell_width, cell_height, is_header)
                
                for paragraph in cell.text_frame.paragraphs:
                    paragraph.alignment = PP_ALIGN.CENTER
                    paragraph.font.size = font_size
                    
                    if is_header:
                        paragraph.font.bold = True
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(46, 92, 138)
                        for run in paragraph.runs:
                            run.font.color.rgb = RGBColor(255, 255, 255)
                    else:
                        if i % 2 == 1:
                            cell.fill.solid()
                            cell.fill.fore_color.rgb = RGBColor(240, 240, 240)
        
        print(f"    ✅ 自适应表格创建成功: {rows}行 x {cols}列")
        print("    ✓ 表格页填充成功")
        
    except Exception as e:
        print(f"    ✗ 表格页填充失败: {e}")
        import traceback
        traceback.print_exc()

# ==================== 版式4：左文右图 ====================
def fill_left_text_right_image_slide(slide, data):
    """填充左文右图页 (版式4)
    占位符: idx=0(标题), idx=10,11,12(页脚), idx=13(文字), idx=14(图片)
    """
    try:
        # 标题 (idx=0)
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:
                placeholder.text = data.get("title", "")
                break
        
        # 文字 (idx=13)
        if "text" in data:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 13:
                    placeholder.text = data["text"]
                    break
        
        # 图片 (idx=14) - 简化处理，仅记录
        if "image_path" in data:
            print(f"    ⚠️ 图片处理待实现: {data['image_path']}")
        
        print("    ✓ 左文右图页填充成功")
    except Exception as e:
        print(f"    ✗ 左文右图页填充失败: {e}")

# ==================== 版式5：左图右文 ====================
def fill_left_image_right_text_slide(slide, data):
    """填充左图右文页 (版式5)
    占位符: idx=0(标题), idx=10,11,12(页脚), idx=13(图片), idx=14(文字)
    """
    try:
        # 标题 (idx=0)
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:
                placeholder.text = data.get("title", "")
                break
        
        # 图片 (idx=13)
        if "image_path" in data:
            print(f"    ⚠️ 图片处理待实现: {data['image_path']}")
        
        # 文字 (idx=14)
        if "text" in data:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 14:
                    placeholder.text = data["text"]
                    break
        
        print("    ✓ 左图右文页填充成功")
    except Exception as e:
        print(f"    ✗ 左图右文页填充失败: {e}")

# ==================== 版式6：三列页 ====================
def fill_three_column_slide(slide, data):
    """填充三列页 (版式6)
    占位符: idx=0(标题), idx=10,11,12(页脚), idx=13-18(三列内容)
    """
    try:
        # 标题 (idx=0)
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:
                placeholder.text = data.get("title", "")
                break
        
        # 列标题和内容
        col_titles = [13, 15, 17]  # 三列标题
        col_contents = [14, 16, 18]  # 三列内容
        
        for i in range(3):
            # 列标题
            title_key = f"col{i+1}_title"
            if title_key in data:
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == col_titles[i]:
                        placeholder.text = data[title_key]
                        break
            
            # 列内容
            items_key = f"col{i+1}_items"
            if items_key in data:
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == col_contents[i]:
                        tf = placeholder.text_frame
                        tf.clear()
                        for j, item in enumerate(data[items_key]):
                            if j == 0:
                                tf.paragraphs[0].text = f"• {item}"
                            else:
                                p = tf.add_paragraph()
                                p.text = f"• {item}"
                        break
        
        print("    ✓ 三列页填充成功")
    except Exception as e:
        print(f"    ✗ 三列页填充失败: {e}")

# ==================== 版式7：四列页 ====================
def fill_four_column_slide(slide, data):
    """填充四列页 (版式7)
    占位符: idx=0(标题), idx=10,11,12(页脚), idx=13-20(四列内容)
    """
    try:
        # 标题 (idx=0)
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:
                placeholder.text = data.get("title", "")
                break
        
        # 四列标题和内容
        col_titles = [13, 15, 17, 19]  # 四列标题
        col_contents = [14, 16, 18, 20]  # 四列内容
        
        for i in range(4):
            # 列标题
            title_key = f"col{i+1}_title"
            if title_key in data:
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == col_titles[i]:
                        placeholder.text = data[title_key]
                        break
            
            # 列内容
            items_key = f"col{i+1}_items"
            if items_key in data:
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == col_contents[i]:
                        tf = placeholder.text_frame
                        tf.clear()
                        for j, item in enumerate(data[items_key]):
                            if j == 0:
                                tf.paragraphs[0].text = f"• {item}"
                            else:
                                p = tf.add_paragraph()
                                p.text = f"• {item}"
                        break
        
        print("    ✓ 四列页填充成功")
    except Exception as e:
        print(f"    ✗ 四列页填充失败: {e}")

# ==================== 版式8：总起页 ====================
def fill_overview_slide(slide, data):
    """填充总起页 (版式8)
    占位符: idx=0(标题), idx=10,11,12(页脚), idx=13-18(数字和要点)
    """
    try:
        # 标题 (idx=0)
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.idx == 0:
                placeholder.text = data.get("title", "内容概览")
                for paragraph in placeholder.text_frame.paragraphs:
                    paragraph.font.size = Pt(32)
                    paragraph.font.bold = True
                    paragraph.alignment = PP_ALIGN.CENTER
                break
        
        # 数字占位符 (13,14,15)
        numbers = ["01", "02", "03"]
        number_indices = [13, 14, 15]
        for idx, num in zip(number_indices, numbers):
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == idx:
                    placeholder.text = num
                    for paragraph in placeholder.text_frame.paragraphs:
                        paragraph.font.size = Pt(20)
                        paragraph.font.bold = True
                        paragraph.alignment = PP_ALIGN.CENTER
                    break
        
        # 三个要点 (16,17,18)
        point_keys = ["point1", "point2", "point3"]
        point_indices = [16, 17, 18]
        
        for i, (key, idx) in enumerate(zip(point_keys, point_indices)):
            if key in data:
                for placeholder in slide.placeholders:
                    if placeholder.placeholder_format.idx == idx:
                        placeholder.text = data[key]
                        for paragraph in placeholder.text_frame.paragraphs:
                            paragraph.font.size = Pt(24)
                            paragraph.font.bold = True
                            paragraph.alignment = PP_ALIGN.CENTER
                        break
        
        print("    ✓ 总起页填充成功")
    except Exception as e:
        print(f"    ✗ 总起页填充失败: {e}")

# ==================== 版式9：标题加多个项目 ====================
def fill_title_content_slide(slide, data):
    """填充标题加内容页 (版式9)
    占位符: idx=10,11,12(页脚), idx=13(标题), idx=14(内容)
    """
    try:
        # 标题 (idx=13)
        if "title" in data:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 13:
                    placeholder.text = data["title"]
                    for paragraph in placeholder.text_frame.paragraphs:
                        paragraph.font.size = Pt(32)
                        paragraph.font.bold = True
                        paragraph.alignment = PP_ALIGN.CENTER
                    print(f"    ✅ 标题填充: {data['title']}")
                    break
        
        # 内容 (idx=14) - 多个项目
        if "items" in data:
            for placeholder in slide.placeholders:
                if placeholder.placeholder_format.idx == 14:
                    tf = placeholder.text_frame
                    tf.clear()
                    
                    items = data["items"]
                    for i, item in enumerate(items):
                        if i == 0:
                            tf.paragraphs[0].text = f"• {item}"
                            tf.paragraphs[0].font.size = Pt(24)
                            tf.paragraphs[0].alignment = PP_ALIGN.LEFT
                            tf.paragraphs[0].space_after = Pt(10)
                        else:
                            p = tf.add_paragraph()
                            p.text = f"• {item}"
                            p.font.size = Pt(24)
                            p.alignment = PP_ALIGN.LEFT
                            p.space_after = Pt(10)
                    
                    print(f"    ✅ 内容填充: {len(items)}个项目")
                    break
        
        print("    ✓ 标题内容页填充成功")
        
    except Exception as e:
        print(f"    ✗ 标题内容页填充失败: {e}")

# ==================== 类型映射 ====================
TYPE_TO_LAYOUT = {
    "compare": 0,        # 比较页
    "end": 1,            # 末尾幻灯片
    "title": 2,          # 自定义标题页
    "table": 3,          # 表格页
    "left_text_right_image": 4,  # 左文右图
    "left_image_right_text": 5,  # 左图右文
    "three_column": 6,   # 三列页
    "four_column": 7,    # 四列页
    "overview": 8,       # 总起页
    "title_content": 9   # 标题加内容页
}

TYPE_TO_HANDLER = {
    "compare": fill_compare_slide,
    "end": fill_end_slide,
    "title": fill_my_title_slide,
    "table": fill_table_slide,
    "left_text_right_image": fill_left_text_right_image_slide,
    "left_image_right_text": fill_left_image_right_text_slide,
    "three_column": fill_three_column_slide,
    "four_column": fill_four_column_slide,
    "overview": fill_overview_slide,
    "title_content": fill_title_content_slide
}

class PPTRenderer:
    """PPT渲染器类，用于FastAPI后台任务"""
    
    def __init__(self, template_path=None):
        """初始化渲染器，可以指定模板路径"""
        self.template_path = template_path or "C:/Users/lenovo/EduForge/backend/services/main.pptx"
    
    def render(self, dsl_data, output_path):
        """
        根据DSL数据生成PPT
        
        参数:
            dsl_data: PPTDocument 对象或字典
            output_path: 输出文件路径
        """
        # 检查模板文件
        if not os.path.exists(self.template_path):
            print(f"❌ 找不到模板文件: {self.template_path}")
            return False
        
        # 复制模板
        print(f"✅ 找到模板: {self.template_path}")
        print(f"📂 输出路径: {output_path}")
        shutil.copy2(self.template_path, output_path)
        prs = Presentation(output_path)
        
        # 将Pydantic模型转换为字典（如果是对象的话）
        if hasattr(dsl_data, 'dict'):
            dsl_dict = dsl_data.dict()
            print("✅ PPTDocument对象转换成功")
        else:
            dsl_dict = dsl_data
        
        # 获取统一背景设置
        global_background = dsl_dict.get("background")
        
        # 设置母版背景
        if global_background:
            print(f"\n🎨 设置母版背景...")
            set_master_background(prs, global_background)
        else:
            print(f"\n🎨 使用模板原有背景")
        
        # 清空示例页
        while len(prs.slides) > 0:
            prs.slides._sldIdLst.remove(prs.slides._sldIdLst[-1])
        
        # 获取页面数据 - 统一从 dsl_dict 获取
        pages_list = dsl_dict.get("pages", [])
        
        # 如果 pages_list 是空的，可能是因为原始对象是 PPTDocument
        if not pages_list and hasattr(dsl_data, 'pages'):
            # 手动构建 pages_list
            pages_list = []
            for page in dsl_data.pages:
                page_dict = {
                    "page_type": page.page_type,  # 保持 page_type，不改为 type
                    "title": page.content.title,
                }
                
                # 添加其他字段
                if page.content.subtitle:
                    page_dict["subtitle"] = page.content.subtitle
                if page.content.items:
                    page_dict["items"] = page.content.items
                if page.content.left_title:
                    page_dict["left_title"] = page.content.left_title
                if page.content.left_items:
                    page_dict["left_items"] = page.content.left_items
                if page.content.right_title:
                    page_dict["right_title"] = page.content.right_title
                if page.content.right_items:
                    page_dict["right_items"] = page.content.right_items
                if page.content.table_data:
                    page_dict["table_data"] = page.content.table_data
                if page.content.point1:
                    page_dict["point1"] = page.content.point1
                if page.content.point2:
                    page_dict["point2"] = page.content.point2
                if page.content.point3:
                    page_dict["point3"] = page.content.point3
                
                pages_list.append(page_dict)
        
        print(f"\n📊 开始生成 {len(pages_list)} 页PPT...")
        
        for page_idx, page_data in enumerate(pages_list):
            # 统一使用 page_type 字段
            page_type = page_data.get("page_type", "content")
            
            # 根据type获取版式索引和处理函数
            layout_idx = TYPE_TO_LAYOUT.get(page_type)
            handler = TYPE_TO_HANDLER.get(page_type)
            
            if layout_idx is None or handler is None:
                print(f"\n⚠️ 未知页面类型: {page_type}，跳过")
                continue
            
            print(f"\n📄 生成第{page_idx+1}页: 类型={page_type}, 版式={layout_idx}")
            
            try:
                # 添加幻灯片
                slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
                
                # 填充内容
                handler(slide, page_data)
                    
            except Exception as e:
                print(f"    ❌ 生成失败: {e}")
                import traceback
                traceback.print_exc()
                continue
            
        # 保存
        prs.save(output_path)
        print(f"\n✅ 成功生成 {len(pages_list)} 页PPT: {output_path}")
        
        return True